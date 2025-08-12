import streamlit as st
import pandas as pd
import plotly.express as px
import io
import json
import google.generativeai as genai
import time
import os
import tempfile
import re
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Page Configuration ---
st.set_page_config(
    page_title="Insight Engine: From CSV to Strategy",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- Session State Initialization ---
if 'analysis_state' not in st.session_state:
    st.session_state.analysis_state = 'initial'
if 'all_report_text' not in st.session_state:
    st.session_state.all_report_text = ""
if 'chart_image_paths' not in st.session_state:
    st.session_state.chart_image_paths = []
if 'df' not in st.session_state:
    st.session_state.df = None
if 'insights_json' not in st.session_state:
    st.session_state.insights_json = []
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'uploaded_file' not in st.session_state:
    st.session_state.uploaded_file = None


# --- Data Sanitization Function ---
def sanitize_columns(df):
    """
    Cleans DataFrame column names to be AI-friendly.
    - Converts to lowercase
    - Replaces spaces and special characters with underscores
    - Removes trailing/leading underscores
    """
    clean_cols = {}
    for col in df.columns:
        # 1. Convert to lowercase
        new_col = col.lower()
        # 2. Replace any character that is not a letter, number, or underscore with an underscore
        new_col = re.sub(r'[^a-zA-Z0-9_]', '_', new_col)
        # 3. Replace multiple consecutive underscores with a single one
        new_col = re.sub(r'_+', '_', new_col)
        # 4. Remove leading or trailing underscores that might result from the above steps
        new_col = new_col.strip('_')
        clean_cols[col] = new_col

    df = df.rename(columns=clean_cols)
    return df


# --- Functions for PDF/PPT Generation and Markdown Parsing ---
def parse_markdown_to_reportlab(text):
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='H1', parent=styles['h1'], fontName='Helvetica-Bold', fontSize=20, leading=24,
                              spaceAfter=14))
    styles.add(ParagraphStyle(name='H2', parent=styles['h2'], fontName='Helvetica-Bold', fontSize=16, leading=20,
                              spaceAfter=10))
    styles.add(ParagraphStyle(name='H3', parent=styles['h3'], fontName='Helvetica-BoldOblique', fontSize=12, leading=14,
                              spaceAfter=8))
    styles.add(ParagraphStyle(name='Body', parent=styles['Normal'], fontName='Helvetica', fontSize=11, leading=14,
                              spaceAfter=6))
    styles.add(ParagraphStyle(name='List', parent=styles['Body'], leftIndent=20, spaceAfter=4))
    story = []
    for line in text.split('\n'):
        stripped_line = line.strip()
        if stripped_line.startswith('```') or stripped_line.startswith('|'):
            continue  # Skip code blocks and markdown tables for PDF
        if stripped_line.startswith('# '):
            story.append(Paragraph(stripped_line[2:], styles['H1']))
            story.append(PageBreak())
        elif stripped_line.startswith('## '):
            story.append(Paragraph(stripped_line[3:], styles['H2']))
        elif stripped_line.startswith('### '):
            story.append(Paragraph(stripped_line[4:], styles['H3']))
        elif stripped_line.startswith('- ') or stripped_line.startswith('* '):
            bullet_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', stripped_line[2:])
            story.append(Paragraph(f'&bull; {bullet_text}', styles['List']))
        elif stripped_line:
            body_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', stripped_line)
            story.append(Paragraph(body_text, styles['Body']))
        else:
            story.append(Spacer(1, 10))
    return story


def create_pdf_report(report_text, image_paths):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, leftMargin=0.75 * inch, rightMargin=0.75 * inch,
                            topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name='Title', parent=styles['h1'], fontSize=24, alignment=1, spaceAfter=20)
    story = [Paragraph("Automated Strategic Insights Report", title_style), Spacer(1, 0.2 * inch)]
    parsed_story = parse_markdown_to_reportlab(report_text)
    story.extend(parsed_story)
    if image_paths:
        story.append(PageBreak())
        story.append(Paragraph("Visualizations", getSampleStyleSheet()['h1']))
        for img_path in image_paths:
            try:
                img = Image(img_path)
                page_width, page_height = letter
                max_width, max_height = page_width - 1.5 * inch, page_height - 2.5 * inch
                scale = min(max_width / img.drawWidth, max_height / img.drawHeight)
                img.drawWidth, img.drawHeight = img.drawWidth * scale, img.drawHeight * scale
                story.append(img)
                story.append(Spacer(1, 0.2 * inch))
            except Exception as e:
                # This warning can be shown in the Streamlit UI if needed
                print(f"Could not add image {img_path} to PDF. It may have been deleted. Error: {e}")
                continue
    doc.build(story)
    buffer.seek(0)
    return buffer


def create_ppt_report(report_text, image_paths, insights_json):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated Strategic Insights Report"
    slide.placeholders[1].text = f"Generated on: {pd.to_datetime('now', utc=True).strftime('%Y-%m-%d %H:%M %Z')}"

    def add_content_slide(title, content_text):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Clean content for PPT
        cleaned_content = re.sub(r'\|.*\|', '', content_text)  # remove tables
        cleaned_content = re.sub(r'`', '', cleaned_content)  # remove backticks
        cleaned_content = re.sub(r'\*\*(.*?)\*\*', r'\1', cleaned_content)  # remove bold markdown
        p.text = cleaned_content
        p.font.size = Pt(14)

    sections = re.split(r'\n## ', '\n' + report_text)
    for section in sections[1:]:
        if not section.strip(): continue
        parts = section.split('\n', 1)
        title, content = (parts[0].strip(), parts[1].strip() if len(parts) > 1 else "")
        if "visualizations" not in title.lower() and "sanitization" not in title.lower():
            add_content_slide(title, content)

    for insight, img_path in zip(insights_json, image_paths):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
            slide.shapes.title.text = insight.get('title', 'Visualization')
            slide.placeholders[1].insert_picture(img_path)
            tf = slide.placeholders[2].text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "Insight: " + insight.get('description', 'No description provided.')
            p.font.size = Pt(16)
        except Exception as e:
            print(f"Could not add image {img_path} to PPT. It may have been deleted. Error: {e}")
            continue
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# --- Gemini API Call Functions ---
@st.cache_data
def get_model(model_name="gemini-1.5-flash-latest"):
    return genai.GenerativeModel(model_name=model_name)


def call_gemini_api_with_backoff(prompt_text, retries=5, json_schema=None, model_name="gemini-1.5-flash-latest"):
    model = get_model(model_name)
    for i in range(retries):
        try:
            config = genai.types.GenerationConfig(
                response_mime_type="application/json" if json_schema else "text/plain",
                temperature=0.2
            )
            if json_schema:
                config.response_schema = json_schema

            response = model.generate_content([{"text": prompt_text}], generation_config=config)
            if response and response.candidates and response.candidates[0].content.parts:
                return response
            else:
                raise ValueError("Received an empty or invalid response from the API.")
        except Exception as e:
            if i < retries - 1:
                st.warning(f"API call failed, retrying in {2 ** i}s... Error: {e}")
                time.sleep(2 ** i)
            else:
                st.error(f"API call failed after multiple retries. Error: {e}")
                return None
    return None


# --- Plotting Functions ---
def plot_graphs_from_insights(df, insights, show_header=True):
    chart_image_paths = []
    visual_insights_text = ""
    plotted_charts_count = 0

    if show_header:
        visual_insights_text = "\n\n## Key Visualizations\n"

    if not insights:
        st.warning("The AI could not generate specific visualizations for this dataset.")
        return visual_insights_text, chart_image_paths

    for i, insight in enumerate(insights):
        title = insight.get('title', f"Untitled Chart {i + 1}")
        description = insight.get('description', "No detailed insight provided by AI.")

        visual_insights_text += f"\n### {title}\n**Insight:** {description}\n"
        st.subheader(f"📈 {title}")
        st.markdown(f"**Insight:** {description}")

        fig = None
        plot_type, x_col, y_col, color_col = insight.get('type'), insight.get('x_axis'), insight.get(
            'y_axis'), insight.get('color')
        valid_cols = df.columns.tolist()

        is_valid = True
        for col, col_name in [(x_col, 'X-axis'), (y_col, 'Y-axis'), (color_col, 'Color')]:
            if col and col not in valid_cols:
                st.warning(
                    f"Skipping plot '{title}': The {col_name} column '{col}' suggested by the AI does not exist in the dataset. Available columns are: `{valid_cols}`")
                is_valid = False
                break
        if not is_valid:
            continue

        try:
            color_palette_discrete = px.colors.qualitative.G10
            color_palette_continuous = px.colors.sequential.Tealgrn
            if plot_type == 'heatmap':
                numeric_df = df.select_dtypes(include=['number'])
                if numeric_df.shape[1] > 1:
                    fig = px.imshow(numeric_df.corr(numeric_only=True), text_auto=True, title=title,
                                    color_continuous_scale=color_palette_continuous)
                else:
                    st.warning(f"Skipping heatmap '{title}': Not enough numeric columns for a correlation matrix.")
            elif plot_type == 'histogram':
                if x_col:
                    fig = px.histogram(df, x=x_col, color=color_col, title=title,
                                       color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping histogram '{title}': An X-axis column is required.")
            elif plot_type == 'scatter':
                if x_col and y_col:
                    fig = px.scatter(df, x=x_col, y=y_col, color=color_col, title=title,
                                     color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping scatter plot '{title}': This chart type requires both an X and a Y axis.")
            elif plot_type == 'line':
                if x_col and y_col:
                    try:
                        df[x_col] = pd.to_datetime(df[x_col])
                    except (ValueError, TypeError):
                        pass
                    df_sorted = df.sort_values(by=x_col)
                    fig = px.line(df_sorted, x=x_col, y=y_col, color=color_col, title=title,
                                  color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping line chart '{title}': This chart type requires both an X and a Y axis.")
            elif plot_type == 'bar':
                if x_col:
                    if y_col and y_col in df.columns:
                        fig = px.bar(df, x=x_col, y=y_col, color=color_col, title=title,
                                     color_discrete_sequence=color_palette_discrete)
                    else:
                        count_df = df[x_col].value_counts().reset_index();
                        count_df.columns = [x_col, 'count']
                        fig = px.bar(count_df, x=x_col, y='count', title=title,
                                     color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping bar chart '{title}': An X-axis column is required.")
            elif plot_type in ['box', 'box_plot']:
                if y_col:
                    fig = px.box(df, x=x_col, y=y_col, color=color_col, title=title,
                                 color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping box plot '{title}': This chart type requires a Y-axis.")
            elif plot_type == 'pie':
                if x_col and y_col:
                    fig = px.pie(df, names=x_col, values=y_col, title=title,
                                 color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(
                        f"Skipping pie chart '{title}': This chart requires a 'names' column (from x_axis) and a 'values' column (from y_axis).")
            elif plot_type == 'violin':
                if y_col:
                    fig = px.violin(df, y=y_col, x=x_col, color=color_col, box=True, title=title,
                                    color_discrete_sequence=color_palette_discrete)
                else:
                    st.warning(f"Skipping violin plot '{title}': This chart type requires a Y-axis.")

            if fig:
                fig.update_layout(title_x=0.5, font=dict(family="Arial, sans-serif"))
                st.plotly_chart(fig, use_container_width=True)
                plotted_charts_count += 1
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                    fig.write_image(tmp_img.name, scale=2)
                    chart_image_paths.append(tmp_img.name)
        except Exception as e:
            st.warning(f"Could not plot '{title}'. The data might be unsuitable for this chart type. Error: {e}")
            continue

    if plotted_charts_count == 0 and insights:
        st.error(
            "The AI provided visualization suggestions, but none could be plotted successfully. This may be due to data incompatibility with the suggested charts. Please review the warnings above.")
    return visual_insights_text, chart_image_paths


def chat_plot_request(df, prompt_text):
    column_list = df.columns.tolist()
    prompt_visuals = f"""
    You are a skilled data analyst. Based on the user's request, create the JSON for a single, most appropriate visualization.
    User request: "{prompt_text}"
    Available columns: {column_list}.
    Respond with a single JSON object describing the visualization. Use only columns from the provided list.
    The JSON must contain: 'title', 'type', 'x_axis', 'y_axis' (can be null), 'color' (optional), and a 'description'.
    """
    schema_visuals = {"type": "OBJECT", "properties": {"title": {"type": "STRING"}, "type": {"type": "STRING"},
                                                       "x_axis": {"type": "STRING"},
                                                       "y_axis": {"type": "STRING", "nullable": True},
                                                       "color": {"type": "STRING", "nullable": True},
                                                       "description": {"type": "STRING"}}}
    response_visuals = call_gemini_api_with_backoff(prompt_visuals, json_schema=schema_visuals)
    if response_visuals:
        try:
            insights_json_str = response_visuals.candidates[0].content.parts[0].text
            insight = json.loads(insights_json_str)
            return [insight]
        except (json.JSONDecodeError, Exception) as e:
            st.error(f"Chat visualization processing error: {e}")
            return []
    return []


# --- Core Analysis Logic ---
def perform_full_analysis(df):
    st.session_state.analysis_state = 'analyzing'
    all_report_text = ""
    all_summaries = {}

    with st.spinner("Preparing data for analysis..."):
        original_columns = df.columns.tolist()
        df = sanitize_columns(df)
        sanitized_columns = df.columns.tolist()

        column_mapping_report = "## Column Name Sanitization Report\n\nTo ensure AI compatibility, some column names were automatically cleaned. All analysis and charts will use these new, cleaned names.\n\n| Original Name | Cleaned Name |\n|---|---|\n"
        changed_count = 0
        for orig, san in zip(original_columns, sanitized_columns):
            if orig != san:
                column_mapping_report += f"| `{orig}` | `{san}` |\n"
                changed_count += 1

        if changed_count == 0:
            column_mapping_report = "## Column Name Sanitization Report\n\nAll column names were already AI-compatible. No changes were needed.\n"

    all_report_text += column_mapping_report

    with st.spinner("Stage 1/7: Understanding the Business Context..."):
        st.header("Stage 1: Business Context")
        prompt_domain = f"""
        You are a seasoned business consultant. Analyze these sanitized column names: {df.columns.tolist()}.
        In plain, non-technical language, describe:
        1. The likely industry or business domain this data represents.
        2. What a non-technical user should understand about these columns.
        3. The key business questions this data could help answer.
        """
        response = call_gemini_api_with_backoff(prompt_domain)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_summaries['domain'] = analysis_text
        all_report_text += f"\n\n## 1. Business Context\n{analysis_text}"
        st.markdown(analysis_text)

    with st.spinner("Stage 2/7: Assessing Data Quality..."):
        st.header("Stage 2: Data Quality Assessment")
        data_info_str = f"{df.info(buf=io.StringIO())}\n\n{df.describe(include='all').to_string()}"
        prompt_health = f"""
        You are a data quality analyst. Based on the data profile below, summarize the key data quality issues in simple terms.
        Focus on:
        - Missing Information: Are there significant gaps in the data? What is the business impact?
        - Strange Values: Are there any numbers or categories that look like potential errors (outliers)?
        - Readiness for Analysis: Is the data clean and ready to be used for decision-making?
        DO NOT include code. Write a short, clear summary.
        Data Profile:\n{data_info_str}
        """
        response = call_gemini_api_with_backoff(prompt_health)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_summaries['health'] = analysis_text
        all_report_text += f"\n\n## 2. Data Quality Assessment\n{analysis_text}"
        st.markdown(analysis_text)

    with st.spinner("Stage 3/7: Discovering Hidden Relationships..."):
        st.header("Stage 3: Hidden Relationship Discovery")
        prompt_relationships = f"""
        You are a creative data scientist and storyteller. Analyze these columns: {df.columns.tolist()}.
        Identify and explain 3-4 *hidden relationships* or *derived concepts* that reveal deeper business insights.
        For each concept, explain:
        1. The concept in plain English (e.g., "Customer Value Density").
        2. What it's made from (e.g., "This combines customer spending with purchase frequency").
        3. The business insight it provides (e.g., "This helps us see if we have many low-value customers or a few high-value ones").
        **CRITICAL: DO NOT PROVIDE ANY PYTHON CODE OR TECHNICAL JARGON.**
        """
        response = call_gemini_api_with_backoff(prompt_relationships)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_summaries['relationships'] = analysis_text
        all_report_text += f"\n\n## 3. Hidden Relationship Discovery\n{analysis_text}"
        st.markdown(analysis_text)

    with st.spinner("Stage 4/7: Generating Key Visualizations..."):
        st.header("Stage 4: Key Visualizations")
        st.markdown(column_mapping_report)

        prompt_visuals = f"""
        You are an expert data visualizer creating a report for a busy executive.
        Your task is to provide a JSON array of 8-10 diverse and insightful visualizations.

        **CRITICAL INSTRUCTION: You MUST ONLY use column names from this EXACT sanitized list: {sanitized_columns}. Do not invent, shorten, or alter any column names.**

        **Chart Selection Guidance:**
        - For single numerical columns (e.g., 'age', 'price'), use a **'histogram'** or **'box_plot'**.
        - To compare a numerical column against a categorical one (e.g., 'price' vs 'category'), use a **'bar'** chart or **'box_plot'**.
        - To compare two numerical columns (e.g., 'price' vs 'size'), use a **'scatter'** plot.
        - For chronological data (e.g., 'sales' over 'date'), use a **'line'** chart.
        - To show parts of a whole (e.g., market share), use a **'pie'** chart.

        **IT IS CRITICAL that for every visualization, the `description` field contains a unique, data-driven insight in 2-3 sentences.** Do NOT use generic descriptions like 'This chart shows X vs Y.'

        The JSON output MUST be a perfectly formed array of objects. Each object MUST contain: 'title', 'type', 'x_axis', 'y_axis', 'color', and the mandatory 'description'.
        """
        schema_visuals = {"type": "ARRAY", "items": {"type": "OBJECT", "properties": {"title": {"type": "STRING"},
                                                                                      "type": {"type": "STRING"},
                                                                                      "x_axis": {"type": "STRING"},
                                                                                      "y_axis": {"type": "STRING",
                                                                                                 "nullable": True},
                                                                                      "color": {"type": "STRING",
                                                                                                "nullable": True},
                                                                                      "description": {
                                                                                          "type": "STRING"}},
                                                     "required": ["title", "type", "x_axis", "description"]}}

        response = call_gemini_api_with_backoff(prompt_visuals, json_schema=schema_visuals)
        insights_json = []
        if response:
            try:
                insights_json_str = response.candidates[0].content.parts[0].text
                insights_json = json.loads(insights_json_str)
            except (json.JSONDecodeError, IndexError, Exception) as e:
                st.error(f"Error parsing visualization insights from AI: {e}")

        visual_insights_text, chart_image_paths = plot_graphs_from_insights(df, insights_json)
        all_summaries['visuals'] = visual_insights_text
        all_report_text += visual_insights_text  # This text is generated inside the plot function now
        st.session_state.chart_image_paths = chart_image_paths
        st.session_state.insights_json = insights_json

    with st.spinner("Stage 5/7: Assessing Predictive Potential..."):
        st.header("Stage 5: Assessing Predictive Potential")
        prompt_predictive = f"""
        You are a machine learning strategist. Analyze the dataset with columns: {df.columns.tolist()}.
        Write a short, strategic memo in plain language.
        1.  **The Predictive Goal**: What is the single most valuable business outcome we could predict with this data?
        2.  **Likelihood of Success**: How confident are you? (High/Medium/Low)
        3.  **Key Factors**: What are the top 3-4 factors from the data that would be most important for this prediction?
        **DO NOT discuss algorithms or complex technical details.**
        """
        response = call_gemini_api_with_backoff(prompt_predictive)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_summaries['predictive'] = analysis_text
        all_report_text += f"\n\n## 5. Assessing Predictive Potential\n{analysis_text}"
        st.markdown(analysis_text)

    with st.spinner("Stage 6/7: Identifying Key Segments..."):
        st.header("Stage 6: Identifying Key Segments")
        prompt_clustering = f"""
        You are a market research analyst. Your task is to identify potential customer groups (segments) from the data with columns: {df.columns.tolist()}.
        Describe 3-4 potential groups you might find. For each group, give them a descriptive name and explain their likely characteristics in simple terms.
        Example:
        - **Power Users**: High frequency, high value.
        - **Budget Shoppers**: Low frequency, focus on discounts.
        **Frame this as a discovery of natural groupings. DO NOT discuss clustering algorithms.**
        """
        response = call_gemini_api_with_backoff(prompt_clustering)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_summaries['clustering'] = analysis_text
        all_report_text += f"\n\n## 6. Identifying Key Segments\n{analysis_text}"
        st.markdown(analysis_text)

    with st.spinner("Stage 7/7: Compiling Strategic Summary..."):
        st.header("Stage 7: Strategic Summary & Recommendations")
        full_context = "\n".join(
            [f"## {key.replace('_', ' ').title()}\n{value}" for key, value in all_summaries.items() if value])
        prompt_strategic = f"""
        You are a Chief Strategy Officer writing a final report for the board.
        Based on the full analysis provided below, write a concise, high-level summary.

        **Full Analysis Context:**
        {full_context}

        **Your final report MUST contain:**
        1.  **Executive Summary**: One paragraph summarizing the most critical finding.
        2.  **Top 3 Insights**: A bulleted list of the three most important discoveries.
        3.  **Top 3 Recommendations**: A numbered list of the three most impactful actions the business should take next.
        """
        response = call_gemini_api_with_backoff(prompt_strategic)
        analysis_text = response.text if response else "Analysis failed for this stage."
        all_report_text += f"\n\n## 7. Strategic Summary & Recommendations\n{analysis_text}"
        st.markdown(analysis_text)

    st.session_state.all_report_text = all_report_text
    st.session_state.data_summary = f"Context: {all_summaries.get('domain', '')}\nStrategic Summary: {analysis_text}"
    st.session_state.df = df
    st.session_state.analysis_state = 'done'
    st.balloons()


# --- Main Application Flow ---
st.title("CSV to Strategy: The Insight Engine")
st.markdown(
    "Upload your data and receive a complete strategic analysis in minutes. Built for business leaders, not analysts.")

with st.sidebar:
    st.header("1. API Key Setup")
    st.info("A Gemini API key is required. Get yours from Google AI Studio.")
    api_key = st.text_input("Enter your Gemini API Key:", type="password",
                            help="Your key is used for this session only and is not stored.")
    if api_key:
        try:
            genai.configure(api_key=api_key)
            st.success("Gemini API key configured!", icon="✅")
        except Exception as e:
            st.error(f"Invalid API Key: {e}")
    else:
        st.warning("Please enter your API key to enable analysis.")
    st.header("2. Analysis Workflow")
    st.markdown(
        """
        - **Step 1:** Upload a CSV or Excel file.
        - **Step 2:** The AI performs a 7-stage analysis.
        - **Step 3:** Review, download, or chat with your data.
        """
    )
    st.markdown("---")
    st.write("Developed for non-technical users.")

if st.session_state.analysis_state == 'initial':
    uploaded_file = st.file_uploader("Upload a CSV or Excel file", type=["csv", "xlsx"])
    if uploaded_file and api_key:
        st.session_state.uploaded_file = uploaded_file
        st.session_state.analysis_state = 'analyzing'
        st.rerun()
    elif uploaded_file and not api_key:
        st.error("Please enter your Gemini API key in the sidebar to begin.")

elif st.session_state.analysis_state == 'analyzing':
    st.info("🚀 Analysis initiated. Please wait as the AI Insight Engine processes your data...")
    if st.session_state.uploaded_file:
        try:
            file = st.session_state.uploaded_file
            df = pd.read_csv(file) if file.name.endswith('.csv') else pd.read_excel(file)
            perform_full_analysis(df)
            st.rerun()
        except Exception as e:
            st.error(f"An error occurred while processing your file: {e}")
            st.session_state.analysis_state = 'initial'
            st.rerun()

else:  # 'done'
    st.success("✅ Analysis Complete! Review your strategic report below.")

    # Display the full report in expandable sections
    report_sections = re.split(r'\n## ', '\n' + st.session_state.all_report_text)
    for i, section in enumerate(report_sections[1:]):
        if not section.strip(): continue
        parts = section.split('\n', 1)
        title = parts[0].strip()
        content = parts[1].strip() if len(parts) > 1 else ""

        is_visuals = "visualizations" in title.lower()
        is_summary = "summary" in title.lower()
        is_setup = "sanitization" in title.lower()

        with st.expander(f"## {title}", expanded=(is_visuals or is_summary or is_setup)):
            if is_visuals:
                if st.session_state.insights_json and st.session_state.df is not None:
                    # The function itself handles the subheaders and plots
                    plot_graphs_from_insights(st.session_state.df, st.session_state.insights_json, show_header=False)
            else:
                st.markdown(content, unsafe_allow_html=True)

    st.header("⬇️ Download Full Report")
    col1, col2 = st.columns(2)
    with col1:
        if st.session_state.all_report_text:
            pdf_buffer = create_pdf_report(st.session_state.all_report_text, st.session_state.chart_image_paths)
            st.download_button(label="📥 Download as PDF", data=pdf_buffer, file_name="Strategic_Insights_Report.pdf",
                               mime="application/pdf")
    with col2:
        if st.session_state.all_report_text:
            ppt_buffer = create_ppt_report(st.session_state.all_report_text, st.session_state.chart_image_paths,
                                           st.session_state.insights_json)
            st.download_button(label="📥 Download as PPTX", data=ppt_buffer, file_name="Strategic_Insights_Report.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    st.header("💬 Chat with your Data")
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            if "plot_data" in message:
                plot_graphs_from_insights(st.session_state.df, message["plot_data"], show_header=False)
            else:
                st.markdown(message["content"])

    if prompt := st.chat_input("Ask a follow-up question... (e.g., 'Plot a histogram of the age column')"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                plot_keywords = ['plot', 'graph', 'chart', 'visualize', 'show me']
                is_plot_request = any(keyword in prompt.lower() for keyword in plot_keywords)

                if is_plot_request and st.session_state.df is not None:
                    insights = chat_plot_request(st.session_state.df, prompt)
                    if insights:
                        st.session_state.messages.append({"role": "assistant", "plot_data": insights})
                        plot_graphs_from_insights(st.session_state.df, insights, show_header=False)
                    else:
                        response_text = "I'm sorry, I couldn't generate a plot for that. Could you be more specific about the columns?"
                        st.markdown(response_text)
                        st.session_state.messages.append({"role": "assistant", "content": response_text})
                else:
                    chat_prompt = f"""
                    You are a helpful AI data analysis assistant. You have access to a pandas DataFrame and its comprehensive analysis report.
                    Here is the compressed context from the report:
                    {st.session_state.data_summary}
                    The available (sanitized) columns in the dataframe are: {st.session_state.df.columns.tolist()}

                    Based ONLY on this context and the dataset's columns, provide a concise and insightful answer to the user's question. Do not hallucinate. If the answer isn't in the context, say so.
                    User's question: {prompt}
                    """
                    response = call_gemini_api_with_backoff(chat_prompt)
                    if response:
                        response_text = response.text
                        st.markdown(response_text)
                        st.session_state.messages.append({"role": "assistant", "content": response_text})
                    else:
                        response_text = "I'm sorry, I could not process that request at this time."
                        st.markdown(response_text)
                        st.session_state.messages.append({"role": "assistant", "content": response_text})