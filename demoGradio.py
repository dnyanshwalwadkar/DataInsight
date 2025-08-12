import gradio as gr
import pandas as pd
import plotly.express as px
import io
import json
import google.generativeai as genai
import time
import os
import tempfile
import re
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt


# --- Functions for PDF/PPT Generation and Markdown Parsing (Unchanged) ---
def parse_markdown_to_reportlab(text):
    """Converts markdown text to a list of ReportLab Flowables."""
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='H1', parent=styles['h1'], fontName='Helvetica-Bold', fontSize=20, leading=24,
                              spaceAfter=14))
    styles.add(ParagraphStyle(name='H2', parent=styles['h2'], fontName='Helvetica-Bold', fontSize=16, leading=20,
                              spaceAfter=10))
    styles.add(ParagraphStyle(name='H3', parent=styles['h3'], fontName='Helvetica-BoldOblique', fontSize=12, leading=14,
                              spaceAfter=8))
    styles.add(ParagraphStyle(name='Body', parent=styles['Normal'], fontName='Helvetica', fontSize=11, leading=14,
                              spaceAfter=6))
    styles.add(ParagraphStyle(name='List', parent=styles['Body'], leftIndent=20, spaceAfter=4))
    code_style = styles['Code']
    code_style.fontName = 'Courier'
    code_style.backColor = '#F0F0F0'
    code_style.leftIndent = 10
    code_style.rightIndent = 10
    code_style.borderPadding = 5
    code_style.spaceBefore = 6
    code_style.spaceAfter = 6
    story = []
    in_code_block = False
    code_text = ""
    for line in text.split('\n'):
        stripped_line = line.strip()
        if stripped_line.startswith('```'):
            if in_code_block:
                story.append(Paragraph(code_text.replace('\n', '<br/>'), styles['Code']))
                code_text = ""
            in_code_block = not in_code_block
            continue
        if in_code_block:
            code_text += stripped_line + '\n'
            continue
        if stripped_line.startswith('# '):
            story.append(Paragraph(stripped_line[2:], styles['H1']))
            story.append(PageBreak())
        elif stripped_line.startswith('## '):
            story.append(Paragraph(stripped_line[3:], styles['H2']))
        elif stripped_line.startswith('### '):
            story.append(Paragraph(stripped_line[4:], styles['H3']))
        elif stripped_line.startswith('- ') or stripped_line.startswith('* '):
            bullet_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', stripped_line[2:])
            story.append(Paragraph(f'&bull; {bullet_text}', styles['List']))
        elif stripped_line:
            body_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', stripped_line)
            story.append(Paragraph(body_text, styles['Body']))
        else:
            story.append(Spacer(1, 10))
    return story


def create_pdf_report(report_text, image_paths):
    """Generates a PDF report from text and images."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, leftMargin=0.75 * inch, rightMargin=0.75 * inch,
                            topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(name='Title', parent=styles['h1'], fontSize=24, alignment=1, spaceAfter=20)
    story = [Paragraph("Automated Strategic Insights Report", title_style), Spacer(1, 0.2 * inch)]
    parsed_story = parse_markdown_to_reportlab(report_text)
    story.extend(parsed_story)
    if image_paths:
        story.append(PageBreak())
        story.append(Paragraph("Visualizations", getSampleStyleSheet()['h1']))
        for img_path in image_paths:
            try:
                img = Image(img_path)
                page_width, page_height = letter
                max_width, max_height = page_width - 1.5 * inch, page_height - 2.5 * inch
                scale = min(max_width / img.drawWidth, max_height / img.drawHeight)
                img.drawWidth, img.drawHeight = img.drawWidth * scale, img.drawHeight * scale
                story.append(img)
                story.append(Spacer(1, 0.2 * inch))
            except Exception as e:
                print(f"Warning: Could not add image {img_path} to PDF. Error: {e}")
                continue
    doc.build(story)
    buffer.seek(0)
    return buffer


def create_ppt_report(report_text, image_paths, insights_json):
    """Generates a PPTX report from text, images, and insights."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automated Strategic Insights Report"
    slide.placeholders[1].text = f"Generated on: {pd.to_datetime('now', utc=True).strftime('%Y-%m-%d %H:%M')}"

    def add_content_slide(title, content_text):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content_text
        p.font.size = Pt(14)

    sections = re.split(r'\n## ', '\n' + report_text)
    for section in sections[1:]:
        if not section.strip(): continue
        parts = section.split('\n', 1)
        title, content = (parts[0].strip(), parts[1].strip() if len(parts) > 1 else "")
        if "visualizations" not in title.lower():
            add_content_slide(title, content)

    for insight, img_path in zip(insights_json, image_paths):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
            slide.shapes.title.text = insight.get('title', 'Visualization')
            slide.placeholders[1].insert_picture(img_path)
            tf = slide.placeholders[2].text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "Insight: " + insight.get('description', 'No description provided.')
            p.font.size = Pt(16)
        except Exception as e:
            print(f"Warning: Could not add image {img_path} to PPT. Error: {e}")
            continue
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# --- Gemini API Call Functions ---
def get_model(model_name="gemini-1.5-flash"):
    """Initializes and returns the Gemini GenerativeModel."""
    return genai.GenerativeModel(model_name=model_name)


def call_gemini_api_with_backoff(prompt_text, retries=5, json_schema=None, model_name="gemini-1.5-flash"):
    """Calls the Gemini API with exponential backoff for retries."""
    model = get_model(model_name)
    for i in range(retries):
        try:
            config = genai.types.GenerationConfig(
                response_mime_type="application/json" if json_schema else "text/plain",
                response_schema=json_schema,
                temperature=0.2
            )
            response = model.generate_content([{"text": prompt_text}], generation_config=config)
            if response and response.candidates and response.candidates[0].content.parts:
                return response
            else:
                raise ValueError("Received an empty or invalid response from the API.")
        except Exception as e:
            if i < retries - 1:
                print(f"API call failed, retrying in {2 ** i}s... Error: {e}")
                time.sleep(2 ** i)
            else:
                print(f"API call failed after multiple retries. Error: {e}")
                return None
    return None


# --- Plotting Functions ---
def plot_graphs_from_insights(df, insights):
    """Generates a list of Plotly figures from AI-generated insights."""
    figs = []
    if not insights:
        print("Warning: The AI could not generate specific visualizations.")
        return figs

    for insight in insights:
        title = insight.get('title')
        description = insight.get('description')

        if not all([title, description]) or "no description" in description.lower():
            continue

        fig = None
        plot_type, x_col, y_col, color_col = insight.get('type'), insight.get('x_axis'), insight.get(
            'y_axis'), insight.get('color')
        valid_cols = df.columns.tolist()
        if (x_col and x_col not in valid_cols) or \
                (y_col and y_col not in valid_cols) or \
                (color_col and color_col not in valid_cols):
            print(f"Warning: Skipping plot '{title}' due to invalid column names.")
            continue

        try:
            color_palette_discrete = px.colors.qualitative.G10
            color_palette_continuous = px.colors.sequential.Tealgrn
            if plot_type == 'heatmap':
                numeric_df = df.select_dtypes(include=['number'])
                if not numeric_df.empty:
                    fig = px.imshow(numeric_df.corr(numeric_only=True), text_auto=True, title=title,
                                    color_continuous_scale=color_palette_continuous)
            elif plot_type == 'histogram':
                fig = px.histogram(df, x=x_col, color=color_col, title=title,
                                   color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'scatter' and x_col and y_col:
                fig = px.scatter(df, x=x_col, y=y_col, color=color_col, title=title,
                                 color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'line' and x_col and y_col:
                fig = px.line(df, x=x_col, y=y_col, color=color_col, title=title,
                              color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'bar':
                if y_col and y_col in df.columns:
                    fig = px.bar(df, x=x_col, y=y_col, color=color_col, title=title,
                                 color_discrete_sequence=color_palette_discrete)
                else:
                    count_df = df[x_col].value_counts().reset_index()
                    count_df.columns = [x_col, 'count']
                    fig = px.bar(count_df, x=x_col, y='count', title=title,
                                 color_discrete_sequence=color_palette_discrete)
            elif plot_type in ['box', 'box_plot'] and y_col:
                fig = px.box(df, x=x_col, y=y_col, color=color_col, title=title,
                             color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'pie' and x_col and y_col:
                fig = px.pie(df, names=x_col, values=y_col, title=title, color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'violin' and y_col:
                fig = px.violin(df, y=y_col, x=x_col, color=color_col, box=True, title=title,
                                color_discrete_sequence=color_palette_discrete)
            elif plot_type == 'sunburst':
                path_cols = [c for c in [x_col, y_col, color_col] if c and c in df.columns]
                if len(path_cols) >= 2:
                    fig = px.sunburst(df, path=path_cols, title=title)
            elif plot_type == 'density_heatmap' and x_col and y_col:
                fig = px.density_heatmap(df, x=x_col, y=y_col, title=title,
                                         color_continuous_scale=color_palette_continuous)

            if fig:
                fig.update_layout(title_x=0.5, font=dict(family="Arial, sans-serif"))
                figs.append(fig)
        except Exception as e:
            print(f"Warning: Could not plot '{title}'. Error: {e}")
            continue
    return figs


def chat_plot_request(df, prompt_text):
    """Generates a single plot based on a user's chat prompt."""
    column_list = df.columns.tolist()
    prompt_visuals = f"""
    You are a highly skilled data analyst. Based on the user's request, create the JSON for a single, most appropriate visualization.
    User request: "{prompt_text}"
    Available columns: {column_list}.
    Your task is to respond with a single JSON object describing the visualization. Do not use columns not in the list.
    The JSON must contain: 'title', 'type', 'x_axis', 'y_axis' (can be null), 'color' (optional), and a mandatory 'description' explaining the plot's insight.
    The JSON MUST be perfectly formed.
    """
    schema_visuals = {"type": "OBJECT", "properties": {"title": {"type": "STRING"}, "type": {"type": "STRING"},
                                                       "x_axis": {"type": "STRING"},
                                                       "y_axis": {"type": "STRING", "nullable": True},
                                                       "color": {"type": "STRING", "nullable": True},
                                                       "description": {"type": "STRING"}}}
    response_visuals = call_gemini_api_with_backoff(prompt_visuals, json_schema=schema_visuals)
    if response_visuals:
        insights_json_str = response_visuals.candidates[0].content.parts[0].text
        try:
            insight = json.loads(insights_json_str)
            if 'description' not in insight or not insight['description'].strip():
                insight['description'] = "This visualization explores data patterns based on the specified columns."
            return [insight]
        except json.JSONDecodeError:
            print("Warning: AI response was malformed. Attempting to fix...")
            fixed_json_str = re.sub(r'(:\s*")([^"]*)(\s*})$', r'\1\2"\3', insights_json_str)
            try:
                insight = json.loads(fixed_json_str)
                if 'description' not in insight or not insight['description'].strip():
                    insight['description'] = "This visualization explores data patterns based on the specified columns."
                return [insight]
            except Exception as e2:
                print(
                    f"Error: Chat visualization processing error after attempting fix: {e2}\nFailed JSON: {insights_json_str}")
                return []
    return []


# --- Core Analysis Logic ---
def perform_full_analysis(df, progress=gr.Progress()):
    """Performs the full 7-stage analysis of the dataframe."""
    all_report_text = ""
    all_summaries = {}

    # --- Pre-emptive Data Cleaning ---
    if 'id' in df.columns: df = df.drop(columns=['id'])
    if 'Unnamed: 32' in df.columns: df = df.drop(columns=['Unnamed: 32'])

    # --- STAGE 1: Business Context ---
    progress(0.1, desc="Stage 1/7: Understanding Business Context...")
    prompt_domain = f"Analyze the column names {df.columns.tolist()} and describe the likely business domain, what a non-technical user should know, and key business questions."
    response = call_gemini_api_with_backoff(prompt_domain)
    analysis_text = response.text if response else "Analysis failed."
    all_summaries['domain'] = analysis_text
    all_report_text += f"\n## 1. Business Context\n{analysis_text}"

    # --- STAGE 2: Data Quality Assessment ---
    progress(0.25, desc="Stage 2/7: Assessing Data Quality...")
    data_info_str = f"{df.info(buf=io.StringIO())}\n\n{df.describe(include='all').to_string()}"
    prompt_health = f"Summarize data quality issues from the profile below for a manager. Focus on missing info, strange values, and readiness for analysis. No code.\nData Profile:\n{data_info_str}"
    response = call_gemini_api_with_backoff(prompt_health)
    analysis_text = response.text if response else "Analysis failed."
    all_summaries['health'] = analysis_text
    all_report_text += f"\n\n## 2. Data Quality Assessment\n{analysis_text}"

    # --- STAGE 3: Hidden Relationship Discovery ---
    progress(0.4, desc="Stage 3/7: Discovering Hidden Relationships...")
    prompt_relationship_insights = f"Analyze columns: {df.columns.tolist()}. Identify 3-4 hidden relationships or derived concepts. Explain the concept, what it's made from, and the business insight. No code or jargon."
    response = call_gemini_api_with_backoff(prompt_relationship_insights)
    analysis_text = response.text if response else "Analysis failed."
    all_summaries['relationships'] = analysis_text
    all_report_text += f"\n\n## 3. Hidden Relationship Discovery\n{analysis_text}"

    # --- STAGE 4: Key Visualizations ---
    progress(0.55, desc="Stage 4/7: Generating Key Visualizations...")
    clean_column_list = df.columns.tolist()
    prompt_visuals = f"""
    Provide a JSON array of 8-10 diverse visualizations.
    **CRITICAL: ONLY use column names from this list: {clean_column_list}.**
    **CRITICAL: Every `description` field MUST contain a unique, data-driven insight (2-3 sentences).**
    **FORBIDDEN `description`:** 'No description provided.', 'This chart shows X vs Y.'
    The JSON output MUST be a perfectly formed array of objects. Each object MUST contain: 'title', 'type', 'x_axis', 'y_axis', 'color', and 'description'.
    """
    schema_visuals = {"type": "ARRAY", "items": {"type": "OBJECT",
                                                 "properties": {"title": {"type": "STRING"}, "type": {"type": "STRING"},
                                                                "x_axis": {"type": "STRING"},
                                                                "y_axis": {"type": "STRING", "nullable": True},
                                                                "color": {"type": "STRING", "nullable": True},
                                                                "description": {"type": "STRING"}},
                                                 "required": ["title", "type", "x_axis", "description"]}}
    response = call_gemini_api_with_backoff(prompt_visuals, json_schema=schema_visuals)
    insights_json = []
    if response:
        try:
            insights_json_str = response.candidates[0].content.parts[0].text
            insights_json = json.loads(insights_json_str)
        except (json.JSONDecodeError, IndexError, Exception) as e:
            print(f"Error parsing visualization insights: {e}")

    visual_insights_text = "\n\n## 4. Key Visualizations\n"
    for insight in insights_json:
        visual_insights_text += f"\n### {insight.get('title', '')}\n**Insight:** {insight.get('description', '')}\n"
    all_summaries['visuals'] = visual_insights_text
    all_report_text += visual_insights_text

    # --- STAGE 5: Predictive Potential ---
    progress(0.7, desc="Stage 5/7: Assessing Predictive Potential...")
    prompt_predictive = f"As a strategist, analyze columns: {df.columns.tolist()}. Memo should state: 1. The most valuable predictive goal. 2. Likelihood of success (High/Medium/Low). 3. Key factors for prediction. No technical details."
    response = call_gemini_api_with_backoff(prompt_predictive)
    analysis_text = response.text if response else "Analysis failed."
    all_summaries['predictive'] = analysis_text
    all_report_text += f"\n\n## 5. Assessing Predictive Potential\n{analysis_text}"

    # --- STAGE 6: Customer/Data Segmentation ---
    progress(0.85, desc="Stage 6/7: Identifying Key Segments...")
    prompt_clustering = f"As a market analyst, identify 3-4 potential customer segments from columns: {df.columns.tolist()}. Give each a name and describe their likely characteristics. No algorithms or code."
    response = call_gemini_api_with_backoff(prompt_clustering)
    analysis_text = response.text if response else "Analysis failed."
    all_summaries['clustering'] = analysis_text
    all_report_text += f"\n\n## 6. Identifying Key Segments\n{analysis_text}"

    # --- STAGE 7: Strategic Summary ---
    progress(0.95, desc="Stage 7/7: Compiling Strategic Summary...")
    full_context = "\n".join(
        [f"## {key.replace('_', ' ').title()}\n{value}" for key, value in all_summaries.items() if value])
    prompt_strategic = f"As Chief Strategy Officer, write a final report based on the full analysis below. It must contain: 1. Executive Summary. 2. Top 3 Insights (bulleted). 3. Top 3 Recommendations (numbered).\n\nFull Analysis Context:\n{full_context}"
    response = call_gemini_api_with_backoff(prompt_strategic)
    analysis_text = response.text if response else "Analysis failed."
    all_report_text += f"\n\n## 7. Strategic Summary & Recommendations\n{analysis_text}"

    data_summary = f"Context: {all_summaries.get('domain', '')}\nStrategic Summary: {analysis_text}"

    return all_report_text, insights_json, data_summary


# --- Gradio Interface Logic ---
def start_analysis_pipeline(api_key, uploaded_file, progress=gr.Progress(track_tqdm=True)):
    """Main function to run the analysis when a file is uploaded."""
    if not api_key:
        raise gr.Error("Gemini API Key is required. Please enter it in the sidebar.")
    if not uploaded_file:
        raise gr.Error("Please upload a CSV or Excel file to begin.")

    try:
        genai.configure(api_key=api_key)
    except Exception as e:
        raise gr.Error(f"Invalid API Key: {e}")

    try:
        fname = uploaded_file.name
        df = pd.read_csv(fname) if fname.endswith('.csv') else pd.read_excel(fname)
    except Exception as e:
        raise gr.Error(f"Error reading file: {e}")

    # Run the full analysis
    all_report_text, insights_json, data_summary = perform_full_analysis(df, progress)

    # Generate plots
    figs = plot_graphs_from_insights(df, insights_json)

    # Save chart images for reports
    chart_image_paths = []
    for fig in figs:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
            fig.write_image(tmp_img.name, scale=2)
            chart_image_paths.append(tmp_img.name)

    # Generate report files
    pdf_buffer = create_pdf_report(all_report_text, chart_image_paths)
    ppt_buffer = create_ppt_report(all_report_text, chart_image_paths, insights_json)

    # Split report text for display
    report_sections = re.split(r'\n## ', '\n' + all_report_text)
    section_contents = {f"sec_{i}": "" for i in range(1, 8)}
    for i, section in enumerate(report_sections[1:], 1):
        section_contents[f"sec_{i}"] = f"## {section}"

    # Prepare outputs for Gradio components
    # The number of plots can be less than 10, so we pad with None
    padded_figs = figs + [None] * (10 - len(figs))

    return (
        df, all_report_text, insights_json, chart_image_paths, data_summary,  # State objects
        gr.update(visible=True),  # Show results_accordion
        gr.update(visible=True),  # Show download_row
        gr.update(visible=True),  # Show chat_interface
        gr.update(value=pdf_buffer.getvalue(), visible=True),  # PDF button
        gr.update(value=ppt_buffer.getvalue(), visible=True),  # PPT button
        section_contents["sec_1"], section_contents["sec_2"], section_contents["sec_3"],
        section_contents["sec_4"], section_contents["sec_5"], section_contents["sec_6"], section_contents["sec_7"],
        *padded_figs
    )


def respond_to_chat(message, history, df, data_summary):
    """Handles the chat interaction."""
    history.append([message, None])

    plot_keywords = ['plot', 'graph', 'chart', 'visualize', 'show me']
    is_plot_request = any(keyword in message.lower() for keyword in plot_keywords)

    if is_plot_request and df is not None:
        insights = chat_plot_request(df, message)
        if insights:
            fig = plot_graphs_from_insights(df, insights)
            if fig:
                # Gradio chatbot doesn't directly support plots. We'll describe it.
                response_text = f"Here is the plot you requested: **{insights[0].get('title')}**. {insights[0].get('description')}"
                history[-1][1] = response_text
                # A better implementation might show the plot in a separate gr.Plot component
            else:
                history[-1][1] = "I'm sorry, I couldn't generate a plot for that."
        else:
            history[-1][
                1] = "I'm sorry, I couldn't generate a plot for that. Could you be more specific about the columns?"
    else:
        chat_prompt = f"""
        You are a helpful AI data analysis assistant. Use the provided context to answer the user's question.
        Context: {data_summary}
        User's question: {message}
        """
        response = call_gemini_api_with_backoff(chat_prompt)
        response_text = response.text if response else "I'm sorry, I could not process that request."
        history[-1][1] = response_text

    return history


# --- Gradio UI Definition ---
with gr.Blocks(theme=gr.themes.Soft(), title="CSV to Strategy: The Insight Engine") as demo:
    # State objects to hold data across interactions
    df_state = gr.State()
    report_text_state = gr.State()
    insights_json_state = gr.State()
    chart_paths_state = gr.State()
    data_summary_state = gr.State()

    gr.Markdown("# CSV to Strategy: The Insight Engine")
    gr.Markdown(
        "Upload your data and receive a complete strategic analysis in minutes. Built for business leaders, not analysts.")

    with gr.Row():
        with gr.Column(scale=1, min_width=300):
            gr.Markdown("## **Setup & Workflow**")
            api_key_input = gr.Textbox(label="1. Gemini API Key", type="password", placeholder="Enter your key here...")
            file_input = gr.File(label="2. Upload CSV or Excel File", file_types=[".csv", ".xlsx"])
            start_button = gr.Button("Start Analysis", variant="primary")

            gr.Markdown(
                """
                ---
                **Analysis Workflow:**
                - **Step 1:** Enter your Gemini API key.
                - **Step 2:** Upload a CSV or Excel file.
                - **Step 3:** The AI performs a 7-stage analysis.
                - **Step 4:** Review, download, or chat with your data.
                """
            )

        with gr.Column(scale=3):
            with gr.Accordion("Full Strategic Report", open=False, visible=False) as results_accordion:
                # Sections for the report
                section_outputs = []
                for i in range(1, 8):
                    section_outputs.append(gr.Markdown(f"### {i}"))

                # Placeholders for plots
                plot_outputs = []
                with gr.Tab("Visualizations"):
                    gr.Markdown("Key visualizations identified from your data.")
                    for i in range(10):  # Max 10 plots
                        plot_outputs.append(gr.Plot())

            with gr.Row(visible=False) as download_row:
                pdf_button = gr.DownloadButton(label="📥 Download as PDF", value=None, visible=False)
                ppt_button = gr.DownloadButton(label="📥 Download as PPTX", value=None, visible=False)

            with gr.Column(visible=False) as chat_interface:
                gr.Markdown("## 💬 Chat with your Data")
                chatbot = gr.Chatbot(label="Follow-up Questions")
                chat_input = gr.Textbox(label="Ask a question...",
                                        placeholder="e.g., 'Explain the key segments in more detail'")
                chat_input.submit(
                    respond_to_chat,
                    [chat_input, chatbot, df_state, data_summary_state],
                    [chatbot]
                )

    # --- Event Wiring ---
    start_button.click(
        start_analysis_pipeline,
        inputs=[api_key_input, file_input],
        outputs=[
            df_state, report_text_state, insights_json_state, chart_paths_state, data_summary_state,
            results_accordion, download_row, chat_interface,
            pdf_button, ppt_button,
            *section_outputs,
            *plot_outputs
        ]
    )

if __name__ == "__main__":
    demo.launch(debug=True)
